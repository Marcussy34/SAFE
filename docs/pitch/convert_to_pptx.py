"""Convert SAFE hackathon HTML pitch deck to PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Colors (SAFE Console palette) ---
BG       = RGBColor(0x0a, 0x0c, 0x11)
PANEL    = RGBColor(0x14, 0x17, 0x1f)
PANEL_S  = RGBColor(0x18, 0x1c, 0x25)
ACCENT   = RGBColor(0x14, 0xf1, 0x95)   # Solana green
BLUE     = RGBColor(0x38, 0xbd, 0xf8)
DANGER   = RGBColor(0xf8, 0x71, 0x71)
TEXT     = RGBColor(0xf4, 0xf7, 0xfa)
MUTED    = RGBColor(0xa2, 0xa8, 0xb4)
FAINT    = RGBColor(0x6b, 0x72, 0x80)
AMBER    = RGBColor(0xfb, 0xbf, 0x24)

W = Inches(13.333)   # 16:9 widescreen
H = Inches(7.5)


def prs_init():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, x, y, w, h, *, bold=False, italic=False,
                 size=18, color=TEXT, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.name   = font
    return txBox


def add_panel(slide, x, y, w, h, color=PANEL):
    """Add a filled rectangle (panel background)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def add_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def topbar(slide, left_text, right_text, left_color=ACCENT, right_color=FAINT):
    """Render top bar with kicker on left and subtitle on right."""
    add_text_box(slide, "// " + left_text,
                 Inches(0.7), Inches(0.35), Inches(7), Inches(0.4),
                 bold=True, size=11, color=left_color, font="Courier New")
    add_text_box(slide, right_text,
                 Inches(7.5), Inches(0.35), Inches(5.5), Inches(0.4),
                 size=11, color=FAINT, align=PP_ALIGN.RIGHT)


def footbar(slide, left_text, right_text):
    """Render bottom bar."""
    add_text_box(slide, left_text,
                 Inches(0.7), Inches(6.9), Inches(9), Inches(0.35),
                 size=10, color=FAINT)
    add_text_box(slide, right_text,
                 Inches(10), Inches(6.9), Inches(3), Inches(0.35),
                 size=10, color=FAINT, align=PP_ALIGN.RIGHT)


def badge(slide, text, x, y, color=FAINT, bg=PANEL_S):
    """Small inline badge."""
    bw = Inches(len(text) * 0.085 + 0.3)
    bh = Inches(0.28)
    p = add_panel(slide, x, y, bw, bh, color=bg)
    add_text_box(slide, text, x + Inches(0.1), y + Inches(0.03),
                 bw - Inches(0.1), bh, size=10, color=color, bold=True)
    return bw


def metric_panel(slide, x, y, w, h, metric, label, bg=PANEL, metric_color=ACCENT):
    add_panel(slide, x, y, w, h, color=bg)
    add_text_box(slide, metric, x + Inches(0.2), y + Inches(0.15),
                 w - Inches(0.3), Inches(0.6), bold=True, size=28, color=metric_color)
    add_text_box(slide, label, x + Inches(0.2), y + Inches(0.75),
                 w - Inches(0.3), Inches(0.5), size=11, color=MUTED)


# ── Slide 1 — Title ───────────────────────────────────────────────
def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # Brand wordmark top-left
    add_text_box(s, "● SAFE//pitch",
                 Inches(0.7), Inches(0.32), Inches(4), Inches(0.4),
                 bold=True, size=12, color=TEXT, font="Courier New")
    add_text_box(s, "2-minute hackathon pitch",
                 Inches(7), Inches(0.32), Inches(6), Inches(0.4),
                 size=11, color=FAINT, align=PP_ALIGN.RIGHT)

    # Kicker
    add_text_box(s, "// SPEND AUTHORIZATION FIREWALL FOR AGENTS",
                 Inches(0.7), Inches(1.3), Inches(11), Inches(0.35),
                 size=11, color=ACCENT, bold=True, font="Courier New")

    # Headline
    add_text_box(s, "Agents can spend.\nSAFE decides if they should.",
                 Inches(0.7), Inches(1.75), Inches(11), Inches(1.9),
                 bold=True, size=42, color=TEXT)

    # Lede
    add_text_box(s,
        "A buyer-side payment firewall that lets AI agents pay for x402 resources "
        "through Solana allowances without blind wallet access.",
                 Inches(0.7), Inches(3.8), Inches(10), Inches(0.9),
                 size=17, color=MUTED)

    # Badges row
    bx = Inches(0.7)
    by = Inches(4.9)
    bx += badge(s, "✓ working devnet demo", bx, by, ACCENT, PANEL_S) + Inches(0.15)
    bx += badge(s, "x402 + Solana allowances", bx, by, BLUE, PANEL_S) + Inches(0.15)
    badge(s, "external-agent API", bx, by, FAINT, PANEL_S)

    footbar(s, "Open with this line.", "01 / 07")
    add_notes(s,
        'Say: "SAFE is a payment firewall for autonomous agents. Wallet caps say how much '
        'an agent can spend. SAFE decides whether this specific payment should happen at all."')
    return s


# ── Slide 2 — Problem ─────────────────────────────────────────────
def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s)
    topbar(s, "The Problem", "Wallet caps are not judgment.")

    add_text_box(s, "A valid transaction can still be the wrong payment.",
                 Inches(0.7), Inches(1.1), Inches(6.5), Inches(1.1),
                 bold=True, size=26, color=TEXT)

    add_text_box(s,
        "Solana x402 has already processed real payment volume. Even tiny wrong-payment "
        "rates become a large safety problem.",
                 Inches(0.7), Inches(2.3), Inches(6.2), Inches(0.8),
                 size=14, color=MUTED)

    # Metric panels left column
    metric_panel(s, Inches(0.7),  Inches(3.25), Inches(2.9), Inches(1.1),
                 "35M+",  "Solana x402 transactions all time.", PANEL, ACCENT)
    metric_panel(s, Inches(3.75), Inches(3.25), Inches(2.9), Inches(1.1),
                 "$10M+", "Solana x402 payment volume all time.", PANEL, ACCENT)
    metric_panel(s, Inches(0.7),  Inches(4.5),  Inches(5.95), Inches(1.0),
                 "35k–350k",
                 "Applying a 0.1%–1% wrong-payment rate to observed Solana x402 volume.",
                 PANEL_S, AMBER)

    add_text_box(s,
        "Modeled from public Solana x402 volume × wrong-payment rates from x402 attack research.",
                 Inches(0.7), Inches(5.65), Inches(6), Inches(0.4),
                 size=9, color=FAINT, italic=True)

    # Right panel — risks
    add_panel(s, Inches(7.2), Inches(1.1), Inches(5.8), Inches(5.3), PANEL)
    risks = [
        "Fake merchant gets paid",
        "Retry loop drains budget",
        "Wrong category slips through",
        "PII metadata leaks out",
        "Recipient mismatch still signs",
    ]
    for i, risk in enumerate(risks):
        add_text_box(s, "✗  " + risk,
                     Inches(7.5), Inches(1.5 + i * 0.82), Inches(5.2), Inches(0.55),
                     size=15, color=DANGER, bold=True)

    footbar(s,
        "Public dashboards count volume. SAFE blocks and records wrong payments.",
        "02 / 07")
    add_notes(s,
        "Spend about 15 seconds. Make the distinction clear: crypto payment validity is not "
        "the same as user-intent validity. The 35k-350k estimate is a conservative model: "
        "0.1%-1% of 35M transactions.")
    return s


# ── Slide 3 — Solution ────────────────────────────────────────────
def slide_solution(prs):
    s = blank_slide(prs)
    fill_bg(s)
    topbar(s, "The Solution", "Pre-signing policy layer")

    add_text_box(s, "SAFE sits between the agent and money.",
                 Inches(0.7), Inches(1.1), Inches(11.5), Inches(0.8),
                 bold=True, size=30, color=TEXT)

    steps = [
        ("01", "x402 challenge",  "Paid API returns a structured payment requirement."),
        ("02", "Normalize",        "Merchant, amount, token, recipient, URL, reason."),
        ("03", "Policy check",     "Trust, category, cap, replay, PII, allowance."),
        ("04", "Settle or block",  "Approved payments settle. Unsafe ones never sign."),
        ("05", "Audit",            "Decision, reason, and tx receipt."),
    ]

    sw = Inches(2.3)
    sh = Inches(2.5)
    gap = Inches(0.12)
    sx = Inches(0.5)
    sy = Inches(2.1)

    for i, (num, label, desc) in enumerate(steps):
        px = sx + i * (sw + gap)
        add_panel(s, px, sy, sw, sh, PANEL)
        add_text_box(s, num,
                     px + Inches(0.15), sy + Inches(0.15), sw - Inches(0.2), Inches(0.55),
                     bold=True, size=28, color=ACCENT, font="Courier New")
        add_text_box(s, label,
                     px + Inches(0.15), sy + Inches(0.75), sw - Inches(0.2), Inches(0.45),
                     bold=True, size=14, color=TEXT)
        add_text_box(s, desc,
                     px + Inches(0.15), sy + Inches(1.25), sw - Inches(0.2), Inches(1.0),
                     size=11, color=MUTED)

    footbar(s, "SAFE is not a wallet. SAFE is the buyer-side payment firewall.", "03 / 07")
    add_notes(s,
        'Say: "The agent does not sign directly. It asks SAFE. SAFE approves, rejects, or '
        'redacts before any settlement is attempted."')
    return s


# ── Slide 4 — Demo ────────────────────────────────────────────────
def slide_demo(prs):
    s = blank_slide(prs)
    fill_bg(s)
    topbar(s, "Live Demo", "Spend about 60 seconds here")

    add_text_box(s, "External agent to devnet settlement.",
                 Inches(0.7), Inches(1.1), Inches(6.2), Inches(0.9),
                 bold=True, size=28, color=TEXT)

    add_text_box(s,
        "The demo proves the full path, not just policy text.",
                 Inches(0.7), Inches(2.1), Inches(6.2), Inches(0.55),
                 size=14, color=MUTED)

    # Badges
    bx = Inches(0.7)
    by = Inches(2.85)
    bx += badge(s, "✓ approved stats",      bx, by, ACCENT, PANEL_S) + Inches(0.15)
    bx += badge(s, "✗ blocked fake merchant", bx, by, DANGER, PANEL_S) + Inches(0.15)
    badge(s, "audit visible", bx, by, BLUE, PANEL_S)

    # Terminal panel
    add_panel(s, Inches(7.1), Inches(1.0), Inches(6.0), Inches(5.4), PANEL)
    # Header bar
    add_panel(s, Inches(7.1), Inches(1.0), Inches(6.0), Inches(0.42), PANEL_S)
    add_text_box(s, "demo script",
                 Inches(7.3), Inches(1.05), Inches(2.5), Inches(0.32),
                 size=11, color=FAINT, font="Courier New")
    add_text_box(s, "localhost:3000",
                 Inches(10.0), Inches(1.05), Inches(2.8), Inches(0.32),
                 size=11, color=FAINT, align=PP_ALIGN.RIGHT, font="Courier New")

    code = (
        "BASE=http://localhost:3000\n\n"
        "pnpm safe doctor\n"
        "pnpm safe pay $BASE/api/x402/stats --dry-run\n"
        "pnpm safe pay $BASE/api/x402/stats\n"
        "pnpm safe pay $BASE/api/x402/fake-merch --dry-run\n"
        "pnpm safe audit"
    )
    add_text_box(s, code,
                 Inches(7.25), Inches(1.55), Inches(5.7), Inches(4.5),
                 size=13, color=ACCENT, font="Courier New")

    footbar(s,
        "Show: decision, tx signature, Explorer URL — then the audit record via `safe audit`.",
        "04 / 07")
    add_notes(s,
        "Switch to terminal or dashboard here. If live network is slow, run dry-run plus "
        "show a previous Explorer URL and audit record.")
    return s


# ── Slide 5 — Product ─────────────────────────────────────────────
def slide_product(prs):
    s = blank_slide(prs)
    fill_bg(s)
    topbar(s, "What Works Today", "Local/private MVP")

    add_text_box(s, "Not a mockup. A working control plane.",
                 Inches(6.8), Inches(1.1), Inches(6.2), Inches(0.9),
                 bold=True, size=26, color=TEXT)

    metric_panel(s, Inches(6.8), Inches(2.1), Inches(2.9), Inches(1.3),
                 "6", "SAFE API routes for agents", PANEL, ACCENT)
    metric_panel(s, Inches(9.85), Inches(2.1), Inches(3.1), Inches(1.3),
                 "3", "integration paths: HTTP, SDK, CLI", PANEL, TEXT)

    badges_row2 = ["preflight", "dry-run", "pay", "audit", "redaction", "replay guard"]
    bx = Inches(6.8)
    by = Inches(3.65)
    for b in badges_row2:
        bw = badge(s, b, bx, by, FAINT, PANEL_S)
        bx += bw + Inches(0.12)
        if bx > Inches(12.5):
            bx = Inches(6.8)
            by += Inches(0.42)

    # Left: placeholder for screenshots
    add_panel(s, Inches(0.5), Inches(1.0), Inches(6.0), Inches(5.4), PANEL_S)
    add_text_box(s, "[SAFE Console\nscreenshots]",
                 Inches(1.5), Inches(3.0), Inches(4), Inches(1.2),
                 size=16, color=FAINT, align=PP_ALIGN.CENTER)

    footbar(s, "Current MVP: devnet-only, local/private SDK and CLI.", "05 / 07")
    add_notes(s,
        "Point out that the product is intentionally bounded: devnet, local facilitator, "
        "local registry. The important part is the end-to-end control path.")
    return s


# ── Slide 6 — Why Now ─────────────────────────────────────────────
def slide_why_now(prs):
    s = blank_slide(prs)
    fill_bg(s)
    topbar(s, "Why Now", "x402 agent payments went 0 → 150M+ in under a year.")

    add_text_box(s,
        "x402 is the agent payment rail.\nAgentic volume is projected to 187× by 2030.",
                 Inches(0.7), Inches(1.1), Inches(11.5), Inches(1.2),
                 bold=True, size=26, color=TEXT)

    # Today panel
    add_panel(s, Inches(0.7), Inches(2.5), Inches(11.9), Inches(0.75), PANEL_S)
    add_text_box(s,
        "x402 today:  150M+ agent payments (119M on Base, 35M on Solana) and ~$600M "
        "annualized volume in under a year — now governed by the Linux Foundation with "
        "Visa, Stripe, Circle, Google & Microsoft behind it.",
                 Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.65),
                 size=12, color=TEXT)

    # 4 metric panels
    cols = [
        ("2025", "+10,000%", "x402 monthly transaction surge when agents went live.", PANEL, TEXT),
        ("2026", "$8B",      "agentic commerce transaction value this year (Juniper).", PANEL, TEXT),
        ("2030", "$1.5T",    "projected agentic commerce value — ~187× in four years (Juniper).", PANEL, ACCENT),
        ("2030", "$3–5T",    "agentic payments worldwide (McKinsey).", PANEL, TEXT),
    ]
    cw = Inches(2.85)
    gap = Inches(0.13)
    for i, (year, val, desc, bg, mc) in enumerate(cols):
        cx = Inches(0.7) + i * (cw + gap)
        cy = Inches(3.45)
        ch = Inches(2.4)
        add_panel(s, cx, cy, cw, ch, bg)
        add_text_box(s, year, cx + Inches(0.15), cy + Inches(0.1),
                     cw - Inches(0.2), Inches(0.3), size=10, color=FAINT, bold=True)
        add_text_box(s, val, cx + Inches(0.15), cy + Inches(0.45),
                     cw - Inches(0.2), Inches(0.75), bold=True, size=30, color=mc)
        add_text_box(s, desc, cx + Inches(0.15), cy + Inches(1.25),
                     cw - Inches(0.2), Inches(1.0), size=11, color=MUTED)

    # Bottom summary
    add_panel(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.7), PANEL_S)
    add_text_box(s,
        "x402 moves the money; SAFE decides whether it should move. "
        "Every agent payment needs a pre-signing safety check — so SAFE's market is agentic payment volume.",
                 Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.6),
                 size=12, color=TEXT)

    footbar(s, "The agent payment rail is compounding — the safety layer has to scale with it.", "06 / 07")
    add_notes(s,
        "Lead with x402-specific traction: 150M+ agent payments (119M on Base, 35M on Solana), "
        "~$600M annualized volume, a 10,000%+ monthly surge in late 2025, Linux Foundation "
        "governance with Visa, Stripe, Circle, Google and Microsoft. "
        "Juniper Research: agentic commerce $8B in 2026 → $1.5T by 2030 (~187x). "
        "McKinsey: $3-5T by 2030. These are third-party market projections.")
    return s


# ── Slide 7 — Vision ──────────────────────────────────────────────
def slide_vision(prs):
    s = blank_slide(prs)
    fill_bg(s)
    topbar(s, "Grand Vision", "Private-first, shared-by-proof")

    add_text_box(s, "From payment firewall to shared trust layer.",
                 Inches(0.7), Inches(1.1), Inches(6.5), Inches(0.9),
                 bold=True, size=30, color=TEXT)

    add_text_box(s,
        "SAFE starts local. Then unknown payments can be verified by private agents, "
        "turned into signed evidence, and reused as sanitized trust signals.",
                 Inches(0.7), Inches(2.1), Inches(6.2), Inches(0.9),
                 size=14, color=MUTED)

    # Badges
    vision_badges = [
        ("private verifier agent", ACCENT),
        ("shared trust database", BLUE),
        ("Solana proof anchors", FAINT),
        ("Filecoin/IPFS evidence", FAINT),
        ("0G later for AI-native DA", FAINT),
    ]
    bx = Inches(0.7)
    by = Inches(3.2)
    for label, color in vision_badges:
        bw = badge(s, label, bx, by, color, PANEL_S)
        bx += bw + Inches(0.15)
        if bx > Inches(6.5):
            bx = Inches(0.7)
            by += Inches(0.42)

    # Right roadmap panel
    add_panel(s, Inches(7.2), Inches(1.1), Inches(5.8), Inches(5.5), PANEL_S)
    roadmap = [
        ("Today:",  "Local SAFE firewall"),
        ("Next:",   "Verifier agents for unknown merchants"),
        ("Then:",   "Shared trust records"),
        ("Later:",  "Decentralized proof and disputes"),
    ]
    for i, (stage, desc) in enumerate(roadmap):
        add_text_box(s, stage,
                     Inches(7.5), Inches(1.55 + i * 0.95), Inches(1.2), Inches(0.45),
                     bold=True, size=14, color=ACCENT)
        add_text_box(s, desc,
                     Inches(8.75), Inches(1.55 + i * 0.95), Inches(4.0), Inches(0.45),
                     size=14, color=TEXT)

    footbar(s,
        "Closing line: agents need more than wallets. They need payment judgment.",
        "07 / 07")
    add_notes(s,
        'End with: "Autonomous agents need more than wallets. They need payment judgment. '
        'SAFE is the firewall before money moves."')
    return s


# ── Main ──────────────────────────────────────────────────────────
def main():
    prs = prs_init()

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_demo(prs)
    slide_product(prs)
    slide_why_now(prs)
    slide_vision(prs)

    out = "/Users/marcus/Projects/SAFE/docs/pitch/hackathon-deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
